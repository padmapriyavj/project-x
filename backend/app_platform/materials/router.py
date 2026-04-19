from io import BytesIO
from typing import Annotated, Optional
import pymupdf as fitz  # PyPI package PyMuPDF — do not install the unrelated `fitz` package
from fastapi import APIRouter, BackgroundTasks, Depends, File, HTTPException, Query, UploadFile, status
from pptx import Presentation
from postgrest.exceptions import APIError

from app_platform.auth.dependencies import get_current_user
from app_platform.storage.client import get_bucket_name, get_s3_client
from database import get_supabase
from models.user import User

from .schemas import MaterialResponse

router = APIRouter(prefix="/api/v1", tags=["materials"])

_MATERIAL_SELECT = "id,course_id,type,filename,s3_key,processing_status,metadata,created_at,updated_at"
_COURSE_SELECT = "id,professor_id"
_LESSON_SELECT = "id,course_id,material_id"

ALLOWED_EXTENSIONS = {"pdf", "ppt", "pptx"}


def _get_file_extension(filename: str) -> str:
    if "." not in filename:
        return ""
    return filename.rsplit(".", 1)[1].lower()


def _course_row_or_404(course_id: int) -> dict:
    sb = get_supabase()
    try:
        res = sb.table("courses").select(_COURSE_SELECT).eq("id", course_id).single().execute()
    except APIError:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Course not found")
    return res.data


def _material_row_or_404(material_id: int) -> dict:
    sb = get_supabase()
    try:
        res = sb.table("materials").select(_MATERIAL_SELECT).eq("id", material_id).single().execute()
    except APIError:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Material not found")
    return res.data


def _user_enrolled(user_id: int, course_id: int) -> bool:
    sb = get_supabase()
    r = (
        sb.table("enrollments")
        .select("id")
        .eq("user_id", user_id)
        .eq("course_id", course_id)
        .limit(1)
        .execute()
    )
    return bool(r.data)


def _extract_text_from_material(material_id: int, s3_key: str, file_type: str) -> None:
    sb = get_supabase()
    try:
        sb.table("materials").update({"processing_status": "processing"}).eq("id", material_id).execute()

        response = get_s3_client().get_object(Bucket=get_bucket_name(), Key=s3_key)
        file_bytes = response["Body"].read()

        if file_type == "pdf":
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pages_text = []
            for page in doc:
                pages_text.append(page.get_text())
            doc.close()
            full_text = "\n".join(pages_text)
            metadata = {"pages": len(pages_text), "full_text": full_text}
        else:
            prs = Presentation(BytesIO(file_bytes))
            slides_text = []
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    slide_texts.append(run.text.strip())
                slides_text.append(" ".join(slide_texts))
            full_text = "\n".join(slides_text)
            metadata = {"slides": len(slides_text), "full_text": full_text}

        sb.table("materials").update({
            "processing_status": "ready",
            "metadata": metadata,
        }).eq("id", material_id).execute()

    except Exception:
        sb.table("materials").update({"processing_status": "failed"}).eq("id", material_id).execute()


@router.post("/courses/{course_id}/materials", response_model=MaterialResponse)
async def upload_material(
    course_id: int,
    background_tasks: BackgroundTasks,
    current_user: Annotated[User, Depends(get_current_user)],
    file: UploadFile = File(...),
    lesson_id: Optional[int] = Query(None),
) -> MaterialResponse:
    if current_user.role != "professor":
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Professor access required")

    course = _course_row_or_404(course_id)
    if course["professor_id"] != current_user.id:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Only the course owner can upload materials")

    ext = _get_file_extension(file.filename or "")
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid file type. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
        )

    file_type = "pdf" if ext == "pdf" else "ppt"
    filename = file.filename or f"upload.{ext}"
    s3_key = f"materials/{course_id}/{filename}"

    file_bytes = await file.read()
    get_s3_client().put_object(Bucket=get_bucket_name(), Key=s3_key, Body=file_bytes)

    sb = get_supabase()
    try:
        ins = sb.table("materials").insert({
            "course_id": course_id,
            "type": file_type,
            "filename": filename,
            "s3_key": s3_key,
            "processing_status": "pending",
            "metadata": {},
        }).execute()
    except APIError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Could not create material",
        ) from e

    if not ins.data:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Material creation failed",
        )

    material = ins.data[0]

    if lesson_id is not None:
        try:
            lesson_res = sb.table("lessons").select(_LESSON_SELECT).eq("id", lesson_id).single().execute()
            if lesson_res.data and lesson_res.data["course_id"] == course_id:
                sb.table("lessons").update({"material_id": material["id"]}).eq("id", lesson_id).execute()
        except APIError:
            pass

    background_tasks.add_task(_extract_text_from_material, material["id"], s3_key, file_type)

    return MaterialResponse.model_validate(material)


@router.get("/courses/{course_id}/materials", response_model=list[MaterialResponse])
def list_materials(
    course_id: int,
    current_user: Annotated[User, Depends(get_current_user)],
) -> list[MaterialResponse]:
    course = _course_row_or_404(course_id)
    if course["professor_id"] != current_user.id and not _user_enrolled(current_user.id, course_id):
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Access denied")

    sb = get_supabase()
    res = sb.table("materials").select(_MATERIAL_SELECT).eq("course_id", course_id).execute()
    return [MaterialResponse.model_validate(row) for row in (res.data or [])]


@router.get("/materials/{material_id}", response_model=MaterialResponse)
def get_material(
    material_id: int,
    current_user: Annotated[User, Depends(get_current_user)],
) -> MaterialResponse:
    material = _material_row_or_404(material_id)
    return MaterialResponse.model_validate(material)


@router.delete("/materials/{material_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_material(
    material_id: int,
    current_user: Annotated[User, Depends(get_current_user)],
) -> None:
    if current_user.role != "professor":
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Professor access required")

    material = _material_row_or_404(material_id)
    course = _course_row_or_404(material["course_id"])

    if course["professor_id"] != current_user.id:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Only the course owner can delete materials")

    try:
        get_s3_client().delete_object(Bucket=get_bucket_name(), Key=material["s3_key"])
    except Exception:
        pass

    sb = get_supabase()
    sb.table("lessons").update({"material_id": None}).eq("material_id", material_id).execute()
    sb.table("materials").delete().eq("id", material_id).execute()
